# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

NAVY_BG = RGBColor(0x0F, 0x17, 0x2A)      # Dark Navy
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)     # Slate Light
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)      # White
PRIMARY_BLUE = RGBColor(0x25, 0x63, 0xEB) # Blue
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)  # Cyan
TEXT_DARK = RGBColor(0x0F, 0x17, 0x2A)    # Slate 900
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)   # Slate 500
BORDER_COLOR = RGBColor(0xCB, 0xD5, 0xE1) # Slate 300

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_header(slide, tag_text, title_text, desc_text=""):
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4), Inches(0.35))
    tf_tag = tag_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = tag_text
    p_tag.font.name = '맑은 고딕'
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = PRIMARY_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.55))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = '맑은 고딕'
    p_title.font.size = Pt(19)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_DARK

    if desc_text:
        desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.35))
        tf_desc = desc_box.text_frame
        p_desc = tf_desc.paragraphs[0]
        p_desc.text = desc_text
        p_desc.font.name = '맑은 고딕'
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_MUTED

# =============================================================
# 슬라이드 1: 표지
# =============================================================
slide_cover = prs.slides.add_slide(blank_layout)
set_slide_background(slide_cover, NAVY_BG)

rect = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.2))
rect.fill.solid()
rect.fill.fore_color.rgb = PRIMARY_BLUE
rect.line.fill.background()

sub_box = slide_cover.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10), Inches(0.5))
tf_sub = sub_box.text_frame
p_sub = tf_sub.paragraphs[0]
p_sub.text = "OVERSEAS MISSION WEB PORTAL SYSTEM"
p_sub.font.name = '맑은 고딕'
p_sub.font.size = Pt(14)
p_sub.font.bold = True
p_sub.font.color.rgb = ACCENT_CYAN

main_box = slide_cover.shapes.add_textbox(Inches(1.2), Inches(2.65), Inches(11), Inches(1.8))
tf_main = main_box.text_frame
tf_main.word_wrap = True
p_main = tf_main.paragraphs[0]
p_main.text = "해외선교부 포털 시스템\n사용자 가이드 & 실무 화면 매뉴얼"
p_main.font.name = '맑은 고딕'
p_main.font.size = Pt(34)
p_main.font.bold = True
p_main.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

desc_box = slide_cover.shapes.add_textbox(Inches(1.2), Inches(4.7), Inches(10), Inches(0.6))
tf_desc = desc_box.text_frame
p_desc = tf_desc.paragraphs[0]
p_desc.text = "실제 UI 화면 스크린샷과 단계별 조작 행동을 담은 사역자 표준 가이드북"
p_desc.font.name = '맑은 고딕'
p_desc.font.size = Pt(13)
p_desc.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

footer_box = slide_cover.shapes.add_textbox(Inches(1.2), Inches(6.2), Inches(10), Inches(0.5))
tf_footer = footer_box.text_frame
p_footer = tf_footer.paragraphs[0]
p_footer.text = "발행 버전 : v2.0 (화면 캡처 수록본)   |   발행 연도 : 2026년   |   대상 : 전도담당자 / 교역자 / 해외선교부"
p_footer.font.name = '맑은 고딕'
p_footer.font.size = Pt(10.5)
p_footer.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

# =============================================================
# 슬라이드 2: 목차
# =============================================================
slide_toc = prs.slides.add_slide(blank_layout)
set_slide_background(slide_toc, LIGHT_BG)
create_header(slide_toc, "TABLE OF CONTENTS", "목 차 (화면별 가이드 구성 안내)", "각 챕터별로 실제 포털 화면 스크린샷과 함께 조작 방법이 안내됩니다.")

toc_data = [
    ("01", "최초 로그인 방법", "로그인 화면 UI, 계정 입력 및 OTP 번호 인증"),
    ("02", "로그인 후 해야하는 것들", "프로필 화면 UI, 비밀번호 변경 및 권한 확인"),
    ("03", "OTP설정 방법", "텔레그램 연동 폼, 봇 검색 및 Chat ID 자동 등록"),
    ("04", "주간보고 작성방법", "주간보고 화면 UI, 5대 부서별 지표 기입 및 제출"),
    ("05", "계획 작성 방법", "전도 계획 탭 UI, 새 블록 추가 및 전략 저장"),
    ("06", "월간보고 작성방법", "월간보고 탭 UI, 활동인원수 입력 및 수정 승인"),
    ("07", "교회별 데이터 확인 방법", "종합 진단 대시보드 UI, Funnel 분석 및 추이 비교"),
    ("부록", "자주 묻는 질문 (FAQ)", "로그인 오류, 텔레그램 미수신 등 문제 해결")
]

for idx, (num, title, desc) in enumerate(toc_data):
    col = idx % 2
    row = idx // 2
    x = Inches(0.8 + col * 5.9)
    y = Inches(1.7 + row * 1.3)
    
    card = slide_toc.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(1.15))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1)
    
    num_box = slide_toc.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2), Inches(0.7), Inches(0.7))
    tf_n = num_box.text_frame
    p_n = tf_n.paragraphs[0]
    p_n.text = num
    p_n.font.name = '맑은 고딕'
    p_n.font.size = Pt(16)
    p_n.font.bold = True
    p_n.font.color.rgb = PRIMARY_BLUE
    
    t_box = slide_toc.shapes.add_textbox(x + Inches(0.85), y + Inches(0.12), Inches(4.5), Inches(0.9))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t1 = tf_t.paragraphs[0]
    p_t1.text = title
    p_t1.font.name = '맑은 고딕'
    p_t1.font.size = Pt(12)
    p_t1.font.bold = True
    p_t1.font.color.rgb = TEXT_DARK
    
    p_t2 = tf_t.add_paragraph()
    p_t2.text = desc
    p_t2.font.name = '맑은 고딕'
    p_t2.font.size = Pt(9.5)
    p_t2.font.color.rgb = TEXT_MUTED

# =============================================================
# 챕터 슬라이드 생성 헬퍼 (UI 스크린샷 전용)
# =============================================================
def add_screen_chapter_slide(tag, title, desc, step_actions, screen_img_filename):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, LIGHT_BG)
    create_header(slide, tag, title, desc)
    
    # 좌측 조작 가이드 카드 (폭 4.8인치)
    card_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.65), Inches(4.6), Inches(5.3))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = BORDER_COLOR
    card_left.line.width = Pt(1)
    
    tb_left = slide.shapes.add_textbox(Inches(0.95), Inches(1.8), Inches(4.3), Inches(5.0))
    tf = tb_left.text_frame
    tf.word_wrap = True
    
    for i, (num_tag, action_title, action_detail) in enumerate(step_actions):
        if i > 0:
            p_head = tf.add_paragraph()
            p_head.space_before = Pt(10)
        else:
            p_head = tf.paragraphs[0]
            
        p_head.text = f"【화면 {num_tag}】 {action_title}"
        p_head.font.name = '맑은 고딕'
        p_head.font.size = Pt(11)
        p_head.font.bold = True
        p_head.font.color.rgb = PRIMARY_BLUE
        
        p_body = tf.add_paragraph()
        p_body.text = action_detail
        p_body.font.name = '맑은 고딕'
        p_body.font.size = Pt(9.2)
        p_body.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
        p_body.space_before = Pt(2)

    # 우측 브라우저 UI 스크린샷 배치 (폭 6.9인치)
    img_path = os.path.join('e:/poject/OverseasWeb/docs_output/images', screen_img_filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(5.6), Inches(1.65), width=Inches(6.9))
    
    return slide

# =============================================================
# 슬라이드 3: 최초 로그인 방법
# =============================================================
add_screen_chapter_slide(
    "CHAPTER 01",
    "1. 최초 로그인 방법",
    "로그인 화면 UI 및 2차 OTP 인증 절차",
    [
        ("[1]번", "아이디 및 초기 임시 비밀번호 입력", "관리자로부터 발급받은 아이디(예: missionary_tokyo)와 초기 임시 비밀번호를 입력창에 넣습니다."),
        ("[2]번", "[로그인] 버튼 클릭", "파란색 [로그인 (Login)] 버튼을 클릭하여 인증을 요청합니다."),
        ("[3]번", "텔레그램 수신 6자리 OTP 입력", "텔레그램으로 즉시 전송된 6자리 번호를 화면 우측 팝업에 입력하고 [인증 완료 및 로그인]을 누릅니다. (유효시간 3분)"),
        ("팁", "텔레그램 자동 로그인 지원", "텔레그램 앱 내 봇 대화창의 [포털 열기] 버튼 클릭 시 무암호 자동 로그인이 지원됩니다.")
    ],
    "screen_01_login.png"
)

# =============================================================
# 슬라이드 4: 로그인 후 해야하는 것들
# =============================================================
add_screen_chapter_slide(
    "CHAPTER 02",
    "2. 로그인 후 해야하는 것들",
    "회원 정보(프로필) 화면 및 필수 비밀번호 변경",
    [
        ("[1]번", "소속 정보 및 배정 범위 점검", "좌측 카드에서 본인의 소속 권한(ROLE_USER)과 데이터 접근 범위(담당 국가 및 교회 목록)가 맞는지 확인합니다."),
        ("[2]번", "새 비밀번호 2회 입력", "우측 카드에서 [현재 비밀번호] 입력 후, 사용할 [새 비밀번호]와 [새 비밀번호 확인]을 동일하게 입력합니다."),
        ("[3]번", "[비밀번호 변경] 클릭", "하단의 [비밀번호 변경 완료] 버튼을 누르면 새 비밀번호가 저장되며 초기 비밀번호 경고가 해제됩니다."),
        ("점검", "상단 네비게이션 확인", "홈/진단, 주간보고, 전도관리(계획/월간) 주요 메뉴 위치를 숙지합니다.")
    ],
    "screen_02_profile_password.png"
)

# =============================================================
# 슬라이드 5: OTP설정 방법
# =============================================================
add_screen_chapter_slide(
    "CHAPTER 03",
    "3. OTP설정 방법 (2단계 보안 인증)",
    "텔레그램 봇 연동 및 Chat ID 등록 설정 화면",
    [
        ("[1]번", "Telegram ID 등록 & 저장", "[회원 정보] 좌측 폼에서 본인의 Telegram ID(@사용자명)를 입력하고 [연동 정보 저장]을 누릅니다."),
        ("[2]번", "텔레그램 봇 대화 시작 (/start)", "텔레그램 앱에서 포털 전용 봇을 검색 후 [시작] 또는 /start를 누르면 Chat ID가 포털에 자동 등록됩니다."),
        ("[3]번", "[테스트 발송] 클릭 및 수신 확인", "우측 하단 [테스트 메시지 발송]을 눌러 텔레그램으로 알림이 정상 도착하는지 확인합니다."),
        ("수동연동", "/myid 명령어 지원", "자동 연동이 안 될 경우 봇에게 /myid 전송 후 알려준 숫자 ID를 직접 입력해 저장할 수 있습니다.")
    ],
    "screen_03_otp_telegram.png"
)

# =============================================================
# 슬라이드 6: 전도담당자로써, 주간보고 작성방법
# =============================================================
add_screen_chapter_slide(
    "CHAPTER 04",
    "4. 전도담당자로써, 주간보고 작성방법",
    "주간보고 작성, 5대 부서 지표 기입, 사진 첨부 및 최종 제출",
    [
        ("[1]번", "교회 및 보고 주차 선택", "상단 드롭다운에서 [담당 교회]와 [작성 주차(예: 8월 3주차)]를 선택합니다."),
        ("[2]번", "5대 부서별 전도 지표 입력", "교역자, 자문회, 장년회, 부녀회, 청년회의 재적, 찾기, 복음방, 가개강, 입교, 탈락수를 표에 직접 기입합니다. (합계 자동 계산)"),
        ("[3]번", "현장 사진 & 특이사항 입력", "우측 사진 영역을 클릭해 현장 사진을 첨부하고 사역 특이사항/기도제목을 작성합니다."),
        ("[4]번", "[최종 제출] 클릭 (마감 준수)", "[임시 저장]으로 중간 보관이 가능하며, 매주 일요일 자정(24:00) 마감 전 [최종 제출]을 완료합니다.")
    ],
    "screen_04_weekly_report.png"
)

# =============================================================
# 슬라이드 7: 계획 작성 방법
# =============================================================
add_screen_chapter_slide(
    "CHAPTER 05",
    "5. 계획 작성 방법 (전도 계획 수립)",
    "전도관리 > 전도 계획(Plan) 탭 작성 및 저장",
    [
        ("[1]번", "[전도 계획 (Plan)] 서브탭 선택", "상단 [전도관리] 메뉴의 3번째 서브탭인 [전도 계획 (Plan)]으로 이동합니다."),
        ("[2]번", "[+ 새 계획 블록 추가] 클릭", "상단 우측의 파란색 버튼을 눌러 새로운 전략 카드 블록을 추가합니다."),
        ("[3]번", "전략 제목 및 실행 내용 작성", "전략 명칭(예: 대학가 집중 캠페인)과 목표 대상, 실행 방안, 기대 효과를 상세히 기입합니다."),
        ("[4]번", "[💾 계획 저장] 클릭", "초록색 [계획 저장 완료]를 클릭하면 서버에 즉시 반영되며 최종 수정자/일시가 기록됩니다.")
    ],
    "screen_05_plan_tab.png"
)

# =============================================================
# 슬라이드 8: 월간보고 작성방법
# =============================================================
add_screen_chapter_slide(
    "CHAPTER 06",
    "6. 월간보고 작성방법",
    "전도관리 > 월간보고 결산, 자동 연산 및 수정 요청",
    [
        ("[1]번", "보고 연도 및 대상 월 선택", "상단 컨트롤 바에서 대상 [연도]와 [월(1월~12월)]을 선택합니다."),
        ("[2]번", "실활동 인원수 & 사역자수 입력", "[수정 모드]를 활성화한 뒤 부서별 [실활동 인원수]와 [전도 사역자(인도자수)]를 입력합니다."),
        ("[3]번", "활동률 & 전년 증감율 자동 계산", "전도재적 대비 활동률(%) 및 전년 동월 대비 증감율이 실시간 자동 산출됩니다."),
        ("[4]번", "과거월 [수정 권한 요청]", "이미 마감된 지난 달 데이터는 [수정 권한 요청] 버튼을 눌러 관리자 승인을 받습니다.")
    ],
    "screen_06_monthly_tab.png"
)

# =============================================================
# 슬라이드 9: 교회별 데이터 확인 방법
# =============================================================
add_screen_chapter_slide(
    "CHAPTER 07",
    "7. 교회별 데이터 확인 방법",
    "종합 진단 대시보드, Funnel 전환 분석, 주차별 추이 그래프",
    [
        ("[1]번", "글로벌 KPI 요약 카드", "전 세계 총 전도재적, 당주차 복음방 개설수, 시온입교 달성 실적을 한눈에 조회합니다."),
        ("[2]번", "단계별 전도 Funnel 퍼널 분석", "[찾기 → 복음방 → 가개강 → 시온입교] 단계별 전환율과 이탈(Drop) 현황을 시각화합니다."),
        ("[3]번", "다중 교회 주차별 비교 추이 그래프", "도쿄, 오사카, 나고야 등 거점 교회의 주차별 성장 곡선을 상호 비교합니다."),
        ("활용", "취합 & 통계표 활용", "부서별 누적 기여도 피벗 테이블 및 엑셀 다운로드 기능을 활용할 수 있습니다.")
    ],
    "screen_07_diagnosis_dashboard.png"
)

# =============================================================
# 슬라이드 10: 부록 FAQ
# =============================================================
slide_faq = prs.slides.add_slide(blank_layout)
set_slide_background(slide_faq, LIGHT_BG)
create_header(slide_faq, "APPENDIX", "부록. 자주 묻는 질문 (FAQ) & 문제 해결", "사역자분들이 자주 겪는 주요 문의 사항과 조치 방법입니다.")

faqs = [
    ("Q. 텔레그램으로 OTP 인증번호가 오지 않습니다.",
     "A. ① [회원 정보]의 Telegram ID가 본인의 실제 @아이디와 일치하는지 확인하세요.\n② 텔레그램 앱에서 포털 전용 봇에게 /start 또는 /myid 메시지를 정상 발송했는지 확인하세요.\n③ 프로필 페이지에서 [테스트 발송]을 눌러 정상 수신 여부를 점검하세요."),
    ("Q. 주간보고 제출 후 오타를 발견했는데 수정이 안 됩니다.",
     "A. 주간보고는 매주 일요일 자정에 자동 마감(Lock)됩니다. 마감 전이라면 언제든 수정 가능하나, 마감된 지난 주차의 경우 상급 관리자에게 문의하여 잠금 해제를 요청하셔야 합니다."),
    ("Q. 담당하지 않는 타 교회의 데이터는 왜 조회가 안 되나요?",
     "A. 시스템 보안 및 데이터 보호 정책상 각 사역자는 본인에게 배정된 국가 및 담당 교회의 데이터만 접근할 수 있습니다. 추가 권한이 필요한 경우 총괄 관리자에게 권한 조정을 요청하세요.")
]

for idx, (q, a) in enumerate(faqs):
    y = Inches(1.7 + idx * 1.7)
    card = slide_faq.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1)
    
    tb = slide_faq.shapes.add_textbox(Inches(1.0), y + Inches(0.12), Inches(11.3), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_q = tf.paragraphs[0]
    p_q.text = q
    p_q.font.name = '맑은 고딕'
    p_q.font.size = Pt(11.5)
    p_q.font.bold = True
    p_q.font.color.rgb = PRIMARY_BLUE
    
    p_a = tf.add_paragraph()
    p_a.text = a
    p_a.font.name = '맑은 고딕'
    p_a.font.size = Pt(9.5)
    p_a.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
    p_a.space_before = Pt(3)

# =============================================================
# 슬라이드 11: 엔딩 슬라이드
# =============================================================
slide_end = prs.slides.add_slide(blank_layout)
set_slide_background(slide_end, NAVY_BG)

end_box = slide_end.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.0))
tf_end = end_box.text_frame
tf_end.word_wrap = True

p_e1 = tf_end.paragraphs[0]
p_e1.text = "감사합니다"
p_e1.font.name = '맑은 고딕'
p_e1.font.size = Pt(36)
p_e1.font.bold = True
p_e1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p_e1.alignment = PP_ALIGN.CENTER

p_e2 = tf_end.add_paragraph()
p_e2.text = "해외선교부 포털 시스템을 통해 전 세계 사역의 승리와 발전을 기원합니다."
p_e2.font.name = '맑은 고딕'
p_e2.font.size = Pt(14)
p_e2.font.color.rgb = ACCENT_CYAN
p_e2.alignment = PP_ALIGN.CENTER
p_e2.space_before = Pt(12)

p_e3 = tf_end.add_paragraph()
p_e3.text = "시스템 이용 및 권한 문의 : 해외선교부 전산관리팀 / 총괄 관리자"
p_e3.font.name = '맑은 고딕'
p_e3.font.size = Pt(11)
p_e3.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
p_e3.alignment = PP_ALIGN.CENTER
p_e3.space_before = Pt(24)

out_pptx = 'e:/poject/OverseasWeb/docs_output/해외선교부_포털_사용자_가이드.pptx'
prs.save(out_pptx)
print(f"Presentation with UI screenshots successfully updated: {out_pptx}")
